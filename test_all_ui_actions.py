import sys
import os
import pytest
from pathlib import Path
from PySide6.QtWidgets import QApplication
from ui import MainWindow, WEB_AVAILABLE
from file_handler import create_viewer_widget

def get_app():
    return QApplication.instance() or QApplication([])

@pytest.fixture
def main_window(tmp_path):
    app = get_app()
    win = MainWindow()
    yield win
    win.close()

def test_mainwindow_instantiation(main_window):
    assert main_window is not None
    assert main_window.tabs is not None

def test_mainwindow_tab_operations(main_window, tmp_path):
    # Test creating new tab
    main_window.new_tab()
    assert main_window.tabs.count() > 0
    
    # Create test file
    test_md = tmp_path / "test.md"
    test_md.write_text("# Hello World\nThis is a test markdown file.")
    
    # Open file
    main_window.open_file(str(test_md))
    assert main_window.tabs.count() > 1
    
    # Close tab
    current_idx = main_window.tabs.currentIndex()
    main_window.close_tab(current_idx)
    
    # Reopen closed tab
    main_window.reopen_closed_tab()

def test_mainwindow_toggle_panels(main_window):
    # Toggle Vault Panel
    main_window.toggle_vault_panel()
    assert main_window.vault_panel is not None
    main_window.toggle_vault_panel()
    
    # Toggle Bookmarks Panel
    main_window.toggle_bookmarks_panel()
    assert main_window.bookmarks_panel is not None
    main_window.toggle_bookmarks_panel()
    
    # Toggle Web Panel
    if WEB_AVAILABLE:
        main_window.toggle_web_panel()
        main_window.toggle_web_panel()

def test_mainwindow_dialog_triggers(main_window, monkeypatch):
    # Prevent exec_ modal blocks during automated test
    monkeypatch.setattr("PySide6.QtWidgets.QDialog.exec", lambda self: 0)
    monkeypatch.setattr("PySide6.QtWidgets.QDialog.exec_", lambda self: 0)
    
    main_window.open_settings()
    main_window.open_feedback_dialog()
    main_window.open_getting_started()

def test_all_menu_actions_and_toolbar_buttons(main_window, monkeypatch):
    # Prevent modal exec blocks
    monkeypatch.setattr("PySide6.QtWidgets.QDialog.exec", lambda self: 0)
    monkeypatch.setattr("PySide6.QtWidgets.QDialog.exec_", lambda self: 0)

    # Test new file creation for all supported extensions
    from editor import EditorTab
    from markdown_renderer import MarkdownViewer
    from docx_viewer import DocxViewer
    from xlsx_viewer import XlsxViewer
    from pptx_viewer import PptxViewer
    from html_viewer import HtmlViewer
    from csv_viewer import CsvViewer

    main_window._create_untitled_tab(".md", MarkdownViewer)
    main_window._create_untitled_tab(".txt", EditorTab)
    main_window._create_untitled_tab(".csv", CsvViewer)
    main_window._create_untitled_tab(".docx", DocxViewer)
    main_window._create_untitled_tab(".xlsx", XlsxViewer)
    main_window._create_untitled_tab(".pptx", PptxViewer)
    main_window._create_untitled_tab(".html", HtmlViewer)

    assert main_window.tabs.count() >= 7

    # Test menu & toolbar slots
    main_window.bookmark_current_tab()
    main_window.toggle_tts_bar()
    main_window.show_find()
    main_window.show_replace()
    main_window.open_vault_search()
    main_window._on_settings_saved()
    main_window.open_quick_switcher()
    main_window._new_session()
    assert main_window.tabs.count() == 1  # Session cleared to welcome/blank tab

def test_all_file_viewers_creation(tmp_path):
    app = get_app()
    
    # 1. Text / Editor
    txt_file = tmp_path / "sample.txt"
    txt_file.write_text("Plain text content")
    v_txt = create_viewer_widget(str(txt_file))
    assert v_txt is not None
    
    # 2. Markdown
    md_file = tmp_path / "sample.md"
    md_file.write_text("# Sample Header\n- List item 1\n- List item 2")
    v_md = create_viewer_widget(str(md_file))
    assert v_md is not None
    
    # 3. CSV
    csv_file = tmp_path / "sample.csv"
    csv_file.write_text("Name,Age,Role\nAlice,30,Developer\nBob,25,Designer")
    v_csv = create_viewer_widget(str(csv_file))
    assert v_csv is not None
    
    # 4. HTML
    html_file = tmp_path / "sample.html"
    html_file.write_text("<html><body><h1>HTML Test</h1></body></html>")
    v_html = create_viewer_widget(str(html_file))
    assert v_html is not None

    # 5. Minimal PDF
    pdf_bytes = b"%PDF-1.4\n1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj 2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj 3 0 obj<</Type/Page/MediaBox[0 0 612 792]>>endobj\nxref\n0 4\n0000000000 65535 f\n0000000009 00000 n\n0000000052 00000 n\n0000000102 00000 n\ntrailer<</Size 4/Root 1 0 R>>\nstartxref\n178\n%%EOF"
    pdf_file = tmp_path / "sample.pdf"
    pdf_file.write_bytes(pdf_bytes)
    v_pdf = create_viewer_widget(str(pdf_file))
    assert v_pdf is not None

    # 6. DOCX
    docx_file = tmp_path / "sample.docx"
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Sample Document", level=1)
        doc.add_paragraph("Hello world paragraph.")
        doc.save(str(docx_file))
    except Exception:
        docx_file.write_bytes(b"%PDF-1.4")
    v_docx = create_viewer_widget(str(docx_file))
    assert v_docx is not None

    # 7. XLSX
    xlsx_file = tmp_path / "sample.xlsx"
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = "Header"
        ws["A2"] = "Value"
        wb.save(str(xlsx_file))
    except Exception:
        xlsx_file.write_bytes(b"%PDF-1.4")
    v_xlsx = create_viewer_widget(str(xlsx_file))
    assert v_xlsx is not None

    # 8. PPTX
    pptx_file = tmp_path / "sample.pptx"
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        prs.save(str(pptx_file))
    except Exception:
        pptx_file.write_bytes(b"%PDF-1.4")
    v_pptx = create_viewer_widget(str(pptx_file))
    assert v_pptx is not None


