"""PowerPoint (.pptx) Viewer component for EleViewer.

Extracts slide text, titles, notes, and bullet points using python-pptx
(with a zipfile XML fallback if python-pptx is unavailable).
Renders slide cards via the bundled Chromium WebEngine (QWebEngineView) for
full CSS fidelity — box shadows, images, themes. Falls back to QTextBrowser
if WebEngine is unavailable.
"""

import html
import os
import re
import zipfile
import xml.etree.ElementTree as ET

try:
    import pptx  # type: ignore
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from PySide6.QtWidgets import (
    QWidget, QVBoxLayout, QHBoxLayout, QToolButton, QLabel,
    QTextBrowser, QLineEdit, QSplitter, QListWidget, QListWidgetItem, QFrame
)
from PySide6.QtCore import Qt, Signal, QSize, QUrl
from PySide6.QtGui import QIntValidator, QKeyEvent, QShortcut, QKeySequence

# Lazy WebEngine availability flag — set inside __init__ to avoid Chromium cold-start tax
_WEB_AVAILABLE: bool | None = None  # None = not yet checked

from icons import icon
from theme import (
    compact_toolbar_stylesheet, ICON_SIZE_COMPACT,
    get_active_palette, get_brand_accent,
)


class PptxViewer(QWidget):
    """Slide-by-slide PowerPoint reader with TTS integration."""

    textChanged = Signal()

    def __init__(self, file_path=None, status_callback=None):
        super().__init__()
        self.file_path = file_path
        self.is_modified = False
        self._status_callback = status_callback
        self._bookmark_callback = None
        self.slides = []  # list of dicts: {"title": str, "content": str, "notes": str}
        self.current_slide = 0
        self.total_slides = 0

        self.setFocusPolicy(Qt.StrongFocus)

        layout = QVBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        layout.setSpacing(0)

        # ── Toolbar ─────────────────────────────────────────────────
        toolbar = QHBoxLayout()
        toolbar.setContentsMargins(6, 4, 6, 4)
        toolbar.setSpacing(4)
        icon_sz = ICON_SIZE_COMPACT
        icon_qsize = QSize(icon_sz, icon_sz)

        def _tb(icon_name, tooltip, slot, text=None):
            btn = QToolButton()
            btn.setIconSize(icon_qsize)
            btn.setIcon(icon(icon_name, size=icon_sz))
            if text:
                btn.setText(f" {text}")
                btn.setToolButtonStyle(Qt.ToolButtonTextBesideIcon)
                btn.setStyleSheet(compact_toolbar_stylesheet() + " QToolButton { font-size: 11px; padding: 2px 6px; }")
            else:
                btn.setStyleSheet(compact_toolbar_stylesheet())
            btn.setToolTip(tooltip)
            btn.setAutoRaise(True)
            btn.clicked.connect(slot)
            return btn

        self.btn_prev = _tb("chevron-left", "Previous Slide", self.prev_slide)
        self.btn_next = _tb("chevron-right", "Next Slide", self.next_slide)
        self.btn_bookmark = _tb("book-open", "Bookmark this slide", self._add_bookmark_here)

        self.slide_input = QLineEdit()
        self.slide_input.setFixedWidth(44)
        self.slide_input.setAlignment(Qt.AlignCenter)
        self.slide_input.setValidator(QIntValidator(1, 9999))
        p = get_active_palette()
        self.slide_input.setStyleSheet(
            f"QLineEdit {{ background:{p['BRAND_PANEL_2']}; color:{p['BRAND_PRIMARY']}; border:1px solid {p['BRAND_BORDER']};"
            f" border-radius:4px; padding:2px 4px; font-weight:bold; font-size:12px; }}"
        )
        self.slide_input.returnPressed.connect(self._jump_to_slide)
        self.lbl_total = QLabel(" / 0")
        self.lbl_total.setStyleSheet(f"color:{p['BRAND_MUTED_FG']}; font-weight:bold; padding:0 6px 0 2px; font-size:12px;")

        toolbar.addWidget(self.btn_prev)
        toolbar.addWidget(self.slide_input)
        toolbar.addWidget(self.lbl_total)
        toolbar.addWidget(self.btn_next)
        toolbar.addStretch()
        toolbar.addWidget(self.btn_bookmark)

        # ── Main Splitter: Slide list + Slide Viewer ────────────────
        self.splitter = QSplitter(Qt.Horizontal)

        self.slide_list = QListWidget()
        self.slide_list.setMaximumWidth(220)
        self.slide_list.setStyleSheet(f"""
            QListWidget {{
                background: {p['BRAND_PANEL']};
                color: {p['BRAND_PRIMARY']};
                border-right: 1px solid {p['BRAND_BORDER']};
                font-size: 12px;
            }}
            QListWidget::item {{ padding: 6px; border-bottom: 1px solid {p['BRAND_BORDER']}; }}
            QListWidget::item:selected {{ background: {get_brand_accent()}; color: {p['BRAND_BACKGROUND']}; font-weight: bold; }}
        """)
        self.slide_list.currentRowChanged.connect(self._on_sidebar_click)

        # ── Search Bar ──────────────────────────────────────────────
        self.search_container = QWidget()
        self.search_container.setStyleSheet(f"background: {p['BRAND_PANEL_2']}; border-bottom: 1px solid {p['BRAND_BORDER']};")
        search_layout = QHBoxLayout(self.search_container)
        search_layout.setContentsMargins(10, 4, 10, 4)
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("Find in presentation (Ctrl+F)...")
        self.search_input.setStyleSheet(f"background: {p['BRAND_PANEL']}; color: {p['BRAND_PRIMARY']}; border: 1px solid {p['BRAND_BORDER']}; border-radius: 4px; padding: 4px;")
        self.search_input.returnPressed.connect(self._perform_search)
        self.btn_search_next = QToolButton()
        self.btn_search_next.setIcon(icon("chevron-down", size=14))
        self.btn_search_next.clicked.connect(self._perform_search)
        search_layout.addWidget(QLabel("Search:"))
        search_layout.addWidget(self.search_input, stretch=1)
        search_layout.addWidget(self.btn_search_next)
        self.search_container.hide()

        # ── Lazy-load WebEngine (Rule 25: never import at module level) ────────
        global _WEB_AVAILABLE
        if _WEB_AVAILABLE is None:
            try:
                from web_panel import _SecureWebView as _SWV, WEB_AVAILABLE as _WA
                _WEB_AVAILABLE = _WA
            except Exception:
                _WEB_AVAILABLE = False

        if _WEB_AVAILABLE:
            from web_panel import _SecureWebView
            self.viewer = _SecureWebView()
            self._use_webengine = True
        else:
            self.viewer = QTextBrowser()
            self.viewer.setStyleSheet(f"""
                QTextBrowser {{
                    background: {p['BRAND_BACKGROUND']};
                    color: {p['BRAND_PRIMARY']};
                    border: none;
                    padding: 0px;
                    font-family: 'Segoe UI', sans-serif;
                }}
            """)
            self.viewer.setOpenExternalLinks(True)
            self._use_webengine = False

        # Sync scroll → sidebar only available in QTextBrowser mode
        if not self._use_webengine:
            self.viewer.verticalScrollBar().valueChanged.connect(self._on_scroll)
        self._is_jumping = False  # prevent circular sync

        right_panel = QWidget()
        right_layout = QVBoxLayout(right_panel)
        right_layout.setContentsMargins(0, 0, 0, 0)
        right_layout.setSpacing(0)
        right_layout.addWidget(self.search_container)
        right_layout.addWidget(self.viewer)

        self.splitter.addWidget(self.slide_list)
        self.splitter.addWidget(right_panel)

        layout.addLayout(toolbar)
        layout.addWidget(self.splitter)
        
        QShortcut(QKeySequence("Ctrl+F"), self).activated.connect(self._show_search)
        
        QShortcut(QKeySequence(Qt.Key_Left), self, context=Qt.WidgetWithChildrenShortcut).activated.connect(self.prev_slide)
        QShortcut(QKeySequence(Qt.Key_Right), self, context=Qt.WidgetWithChildrenShortcut).activated.connect(self.next_slide)

        if file_path:
            self.load_from_path(file_path)

    def _extract_shapes(self, shape):
        """Return list of (kind, value) tuples: ('text', str) or ('image', (bytes, ext))."""
        results = []
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    results.append(("text", text))

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            results.append(("text", text))

        if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                img_bytes = shape.image.blob
                img_ext = shape.image.ext.lower()
                if img_bytes and img_ext:
                    results.append(("image", (img_bytes, img_ext)))
            except Exception as e:
                print(f"[PPTX] Failed to extract image: {e}")

        return results

    def load_from_path(self, file_path):
        self.file_path = file_path
        self.slides = []

        if PPTX_AVAILABLE:
            try:
                prs = pptx.Presentation(file_path)
                for idx, slide in enumerate(prs.slides, start=1):
                    title = f"Slide {idx}"
                    texts = []
                    images = []  # list of (bytes, ext)
                    first_text_seen = False
                    for shape in slide.shapes:
                        for kind, val in self._extract_shapes(shape):
                            if kind == "text":
                                if not first_text_seen and getattr(shape, "is_placeholder", False) and getattr(shape, "placeholder_format", None) is not None:
                                    title = val
                                texts.append(val)
                                first_text_seen = True
                            else:  # image
                                images.append(val)
                    notes = ""
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes = slide.notes_slide.notes_text_frame.text.strip()

                    self.slides.append({
                        "title": title,
                        "content": "\n\n".join(texts),
                        "images": images,
                        "notes": notes,
                    })
            except Exception as e:
                print(f"[PPTX Viewer] python-pptx load failed: {e}")
                self._fallback_zip_parse(file_path)
        else:
            self._fallback_zip_parse(file_path)

        self.total_slides = len(self.slides)
        self.lbl_total.setText(f" / {self.total_slides}")

        self.slide_list.clear()
        for idx, s in enumerate(self.slides, start=1):
            title_text = s["title"][:28] + ("…" if len(s["title"]) > 28 else "")
            item = QListWidgetItem(f"{idx}. {title_text}")
            self.slide_list.addItem(item)

        if self.total_slides > 0:
            self._render_continuous_html()
            self.go_to_slide(0)

    def _fallback_zip_parse(self, file_path):
        """Fallback parser extracting slide text XML from pptx zip container."""
        import base64
        try:
            with zipfile.ZipFile(file_path, "r") as z:
                slide_files = [name for name in z.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
                # Sort numerically to prevent slide10 coming before slide2
                slide_files.sort(key=lambda x: int(x.split("slide")[-1].split(".xml")[0]))
                for idx, sfile in enumerate(slide_files, start=1):
                    xml_content = z.read(sfile)
                    tree = ET.fromstring(xml_content)
                    texts = []
                    for elem in tree.iter():
                        if elem.tag.endswith("}t") and elem.text:
                            t = elem.text.strip()
                            if t:
                                texts.append(t)

                    # Native image extraction via .rels mapping
                    images = []
                    parts = sfile.split("/")
                    rels_path = "/".join(parts[:-1]) + "/_rels/" + parts[-1] + ".rels"
                    if rels_path in z.namelist():
                        rels_tree = ET.fromstring(z.read(rels_path))
                        for rel in rels_tree.iter():
                            if rel.tag.endswith("}Relationship"):
                                target = rel.attrib.get("Target", "")
                                if target.startswith("../media/"):
                                    media_path = "ppt/media/" + target.split("/")[-1]
                                    if media_path in z.namelist():
                                        img_bytes = z.read(media_path)
                                        ext = media_path.split(".")[-1].lower()
                                        images.append((img_bytes, ext))

                    title = f"Slide {idx}"
                    if texts:
                        title = texts[0]
                    self.slides.append({
                        "title": title,
                        "content": "\n\n".join(texts),
                        "images": images,
                        "notes": "",
                    })
        except Exception as e:
            print(f"[PPTX Viewer] Fallback zip parse error: {e}")

    def set_bookmark_callback(self, callback):
        self._bookmark_callback = callback

    def _bookmark_payload(self):
        name = os.path.basename(self.file_path) if self.file_path else "presentation"
        return {
            "page_number": self.current_slide,
            "scroll_position_y": 0.0,
            "label": f"Slide {self.current_slide + 1} in {name}",
        }

    def _add_bookmark_here(self):
        if self._bookmark_callback:
            self._bookmark_callback(self._bookmark_payload())

    def go_to_bookmark(self, page_number=0, scroll_position_y=0.0):
        self.go_to_slide(int(page_number))

    def _show_search(self):
        self.search_container.show()
        self.search_input.setFocus()
        self.search_input.selectAll()

    def _perform_search(self):
        query = self.search_input.text()
        if not query:
            return
        if self._use_webengine:
            # QWebEngineView.findText wraps automatically
            from PySide6.QtWebEngineWidgets import QWebEngineView
            self.viewer.findText(query)
        else:
            # QTextBrowser.find — loop back on no match
            if not self.viewer.find(query):
                self.viewer.moveCursor(self.viewer.textCursor().Start)
                self.viewer.find(query)

    def _render_text_block(self, text, empty_message=""):
        if not text:
            return empty_message
        return html.escape(text).replace("\n", "<br>").replace("\r", "")

    def _render_slide_content(self, raw_content):
        if not raw_content:
            return "<i>(No text content on this slide)</i>"

        fragments = []
        for chunk in re.split(r"(<img[^>]*>)", raw_content):
            if not chunk:
                continue
            if chunk.startswith("<img"):
                fragments.append(chunk)
            else:
                fragments.append(self._render_text_block(chunk))
        return "".join(fragments)

    def _render_continuous_html(self):
        import base64
        p = get_active_palette()
        accent = get_brand_accent()

        if self._use_webengine:
            # WebEngine: inline base64 — no size limit, full CSS support
            html_output = f'<!DOCTYPE html><html><head><meta charset="utf-8">\n<style>\n  body {{ margin: 0; padding: 0; background: {p["BRAND_BACKGROUND"]}; font-family: \'Segoe UI\', sans-serif; }}\n  .slide-card {{ max-width: 800px; margin: 30px auto; padding: 40px; background: {p["BRAND_PANEL"]}; border: 1px solid {p["BRAND_BORDER"]}; border-radius: 8px; box-shadow: 0 4px 16px rgba(0,0,0,0.35); }}\n  .slide-label {{ color: {accent}; font-size: 0.8em; font-weight: bold; text-transform: uppercase; letter-spacing: 1px; margin-bottom: 8px; }}\n  .slide-title {{ color: {p["BRAND_PRIMARY"]}; font-size: 1.5em; margin: 0 0 20px 0; border-bottom: 1px solid {p["BRAND_BORDER"]}; padding-bottom: 10px; }}\n  .slide-body {{ font-size: 1.1em; line-height: 1.7; color: {p["BRAND_PRIMARY"]}; }}\n  .slide-body img {{ max-width: 100%; border-radius: 4px; margin: 10px 0; }}\n  .notes-box {{ margin-top: 30px; padding: 12px; background: {p["BRAND_PANEL_2"]}; border-left: 3px solid {accent}; border-radius: 4px; }}\n  .notes-label {{ color: {p["BRAND_MUTED_FG"]}; font-size: 0.8em; font-weight: bold; margin-bottom: 4px; }}\n  .notes-body {{ font-size: 0.9em; color: {p["BRAND_MUTED_FG"]}; }}\n  .bottom-pad {{ height: 40px; }}\n</style></head><body>'

            for index, slide_data in enumerate(self.slides):
                content_html = self._render_slide_content(slide_data["content"] or "")
                # Append images inline for WebEngine
                for img_bytes, img_ext in slide_data.get("images", []):
                    mime = "image/jpeg" if img_ext == "jpg" else f"image/{img_ext}"
                    b64 = base64.b64encode(img_bytes).decode("utf-8")
                    content_html += f'<img src="data:{mime};base64,{b64}" style="max-width:100%; border-radius:4px; margin:10px 0;"/>'

                title_text = self._render_text_block(slide_data["title"])
                notes_text = self._render_text_block(slide_data.get("notes", ""))

                html_output += f'<div id="slide_{index}" class="slide-card">'
                html_output += f'<div class="slide-label">SLIDE {index + 1} OF {self.total_slides}</div>'
                html_output += f'<h1 class="slide-title">{title_text}</h1>'
                html_output += f'<div class="slide-body">{content_html}</div>'
                if notes_text:
                    html_output += (f'<div class="notes-box">'
                                    f'<div class="notes-label">SPEAKER NOTES</div>'
                                    f'<div class="notes-body">{notes_text}</div></div>')
                html_output += '</div>'

            html_output += '<div class="bottom-pad"></div></body></html>'
            self.viewer.setHtml(html_output, QUrl())

        else:
            # QTextBrowser: register images as document resources to avoid the
            # ~2MB setHtml string limit that silently drops large base64 blobs.
            # Rule 21: do NOT use white-space:pre-wrap on containers with <img>.
            from PySide6.QtGui import QTextDocument, QImage
            from PySide6.QtCore import QUrl as _QUrl

            doc = QTextDocument()
            doc.setDefaultStyleSheet(
                f"body {{ background: {p['BRAND_BACKGROUND']}; font-family: 'Segoe UI', sans-serif; }}"
                f".slide-card {{ max-width: 700px; margin: 24px auto; padding: 30px;"
                f"  background: {p['BRAND_PANEL']}; border: 1px solid {p['BRAND_BORDER']}; border-radius: 6px; }}"
                f".slide-label {{ color: {accent}; font-size: small; font-weight: bold; text-transform: uppercase; }}"
                f".slide-title {{ color: {p['BRAND_PRIMARY']}; font-size: large; margin: 8px 0 16px 0;"
                f"  border-bottom: 1px solid {p['BRAND_BORDER']}; padding-bottom: 8px; }}"
                f".slide-body {{ font-size: medium; line-height: 1.6; color: {p['BRAND_PRIMARY']}; }}"
                f".notes-box {{ margin-top: 20px; padding: 10px; background: {p['BRAND_PANEL_2']};"
                f"  border-left: 3px solid {accent}; border-radius: 4px; }}"
                f".notes-label {{ color: {p['BRAND_MUTED_FG']}; font-size: small; font-weight: bold; }}"
                f".notes-body {{ font-size: small; color: {p['BRAND_MUTED_FG']}; }}"
            )

            img_uris = {}  # (slide_index, img_index) -> URI string
            for s_idx, slide_data in enumerate(self.slides):
                for i_idx, (img_bytes, img_ext) in enumerate(slide_data.get("images", [])):
                    uri = f"pptx-img://{s_idx}/{i_idx}.{img_ext}"
                    img_uris[(s_idx, i_idx)] = uri
                    qimg = QImage()
                    qimg.loadFromData(img_bytes)
                    doc.addResource(QTextDocument.ImageResource, _QUrl(uri), qimg)

            html_output = "<body>"
            for index, slide_data in enumerate(self.slides):
                content_html = self._render_slide_content(slide_data["content"] or "")
                # Append images using registered resource URIs — no base64 in HTML string
                for i_idx in range(len(slide_data.get("images", []))):
                    uri = img_uris[(index, i_idx)]
                    content_html += f'<img src="{uri}" style="max-width:100%; margin:8px 0;"/>'

                title_text = self._render_text_block(slide_data["title"])
                notes_text = self._render_text_block(slide_data.get("notes", ""))

                html_output += f'<a name="slide_{index}"></a>'
                html_output += f'<div class="slide-card">'
                html_output += f'<div class="slide-label">SLIDE {index + 1} OF {self.total_slides}</div>'
                html_output += f'<div class="slide-title">{title_text}</div>'
                html_output += f'<div class="slide-body">{content_html}</div>'
                if notes_text:
                    html_output += (f'<div class="notes-box">'
                                    f'<div class="notes-label">SPEAKER NOTES</div>'
                                    f'<div class="notes-body">{notes_text}</div></div>')
                html_output += '</div>'
            html_output += '</body>'
            doc.setHtml(html_output)
            self.viewer.setDocument(doc)

    def _on_sidebar_click(self, index):
        if not self._is_jumping:
            self.go_to_slide(index)

    def _on_scroll(self, value):
        if self._is_jumping or self.total_slides == 0:
            return
            
        # Estimate which slide is currently in view based on scroll bar percentage
        scrollbar = self.viewer.verticalScrollBar()
        max_val = scrollbar.maximum()
        if max_val <= 0:
            return
            
        percentage = value / max_val
        # Add slight offset so it switches earlier
        estimated_index = min(self.total_slides - 1, int(percentage * self.total_slides + 0.1))
        
        if estimated_index != self.current_slide:
            self.current_slide = estimated_index
            self.slide_input.setText(str(estimated_index + 1))
            
            # Update list without triggering currentRowChanged
            self.slide_list.blockSignals(True)
            self.slide_list.setCurrentRow(estimated_index)
            self.slide_list.blockSignals(False)

    def go_to_slide(self, index):
        if index < 0 or index >= self.total_slides:
            return

        self._is_jumping = True
        self.current_slide = index
        self.slide_input.setText(str(index + 1))

        self.slide_list.blockSignals(True)
        self.slide_list.setCurrentRow(index)
        self.slide_list.blockSignals(False)

        if self._use_webengine:
            # Use JS scrollIntoView for smooth, accurate slide navigation
            js = f"var el = document.getElementById('slide_{index}'); if (el) el.scrollIntoView({{behavior: 'smooth', block: 'start'}});"
            self.viewer.page().runJavaScript(js)
        else:
            self.viewer.scrollToAnchor(f"slide_{index}")

        from PySide6.QtCore import QTimer
        QTimer.singleShot(100, lambda: setattr(self, '_is_jumping', False))

    def prev_slide(self):
        if self.current_slide > 0:
            self.go_to_slide(self.current_slide - 1)

    def next_slide(self):
        if self.current_slide < self.total_slides - 1:
            self.go_to_slide(self.current_slide + 1)

    def _jump_to_slide(self):
        try:
            s = int(self.slide_input.text()) - 1
            self.go_to_slide(s)
        except ValueError:
            pass

    def read_current_page(self, voice_id=None):
        """Return text of current slide for TTS. Playback handled by MainWindow."""
        if 0 <= self.current_slide < self.total_slides:
            slide_data = self.slides[self.current_slide]
            text = f"{slide_data['title']}. {slide_data['content']}"
            return text.strip() if text.strip() else ""
        return ""

    def toPlainText(self):
        """Return all text content across all slides."""
        full_text = []
        for idx, s in enumerate(self.slides, start=1):
            full_text.append(f"--- Slide {idx}: {s['title']} ---\n{s['content']}\n")
        return "\n".join(full_text)

    def setPlainText(self, text):
        pass

    def keyPressEvent(self, event: QKeyEvent):
        key = event.key()
        if key in (Qt.Key_Left, Qt.Key_PageUp):
            self.prev_slide()
            event.accept()
        elif key in (Qt.Key_Right, Qt.Key_PageDown):
            self.next_slide()
            event.accept()
        else:
            super().keyPressEvent(event)
