from pptx import Presentation
from pptx.util import Inches
from PySide6.QtGui import QImage
from PySide6.QtWidgets import QApplication

import pptx_viewer


def test_embedded_images_keep_slide_order(tmp_path, monkeypatch):
    app = QApplication.instance() or QApplication([])
    monkeypatch.setattr(pptx_viewer, "_WEB_AVAILABLE", False)

    image_path = tmp_path / "embedded.png"
    image = QImage(12, 8, QImage.Format.Format_RGB32)
    image.fill(0x2A82DA)
    assert image.save(str(image_path), "PNG")

    presentation_path = tmp_path / "ordered.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    before = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(0.5))
    before.text = "Text before image"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.2), width=Inches(2))
    after = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(0.5))
    after.text = "Text after image"
    presentation.save(str(presentation_path))

    viewer = pptx_viewer.PptxViewer(str(presentation_path))
    try:
        elements = viewer.slides[0]["elements"]
        assert [kind for kind, _ in elements] == ["text", "image", "text"]
        assert elements[0][1] == "Text before image"
        assert elements[2][1] == "Text after image"

        rendered = viewer._render_qtext_elements(elements, ["pptx-img://0/0.png"])
        assert rendered.index("Text before image") < rendered.index("pptx-img://0/0.png")
        assert rendered.index("pptx-img://0/0.png") < rendered.index("Text after image")
    finally:
        viewer.close()
